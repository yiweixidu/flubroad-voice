import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional
from .visualizer import create_neutralization_heatmap

class PPTGenerator:
    def __init__(self, template_path: Optional[str] = None):
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            if template_path:
                print(f"Warning: Template file '{template_path}' not found. Using default blank presentation.")

    def add_title_slide(self, title: str, subtitle: str = ""):
        slide_layout = self.prs.slide_layouts[0]  # title slide
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle

    def add_content_slide(self, title: str, bullets: List[str]):
        slide_layout = self.prs.slide_layouts[1]  # title and content
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0

    def add_table_slide(self, title: str, headers: List[str], rows: List[List[str]]):
        slide_layout = self.prs.slide_layouts[5]  # title only
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        rows_count = len(rows)
        cols_count = len(headers)
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(0.5 * (rows_count + 1))
        table = slide.shapes.add_table(rows_count + 1, cols_count, left, top, width, height).table
        # header row
        for col, header in enumerate(headers):
            table.cell(0, col).text = header
        # data rows
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                table.cell(i+1, j).text = str(cell)

    def add_neutralization_heatmap(self, data: Dict[str, List[float]], title: str = "Neutralization Breadth"):
        """添加热图幻灯片"""
        img_bytes = create_neutralization_heatmap(data, title)
        slide_layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(img_bytes, left, top, height=Inches(5))
        # 添加标题（可选）
        title_box = slide.shapes.add_textbox(left, Inches(0.5), Inches(8), Inches(1))
        title_box.text = title

    def save(self, path: str):
        self.prs.save(path)